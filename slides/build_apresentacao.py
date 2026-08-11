# build_apresentacao.py — gera slides/apresentacao_final.pptx (13 slides,
# ~10 min de fala). Ordem: motivação -> dados e variáveis -> por que o
# furto (outliers por crime) -> correlação -> desenvolvimento do modelo
# (autocorrelação, OLS, espacial global/GWR, contrafactual, mecanismo,
# exemplo São Paulo x Mongaguá, comparação final) -> conclusão. Sem slide
# de metodologia isolado (removido a pedido da orientação). Usa as figuras
# já geradas em outputs/ pelo pipeline (Rscript run_espacial.R).
#
# Uso: python slides/build_apresentacao.py

import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(ROOT, "outputs")
SLIDES_DIR = os.path.join(ROOT, "slides")

NAVY = RGBColor(0x1F, 0x4E, 0x79)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_SOFT = RGBColor(0xE2, 0xEF, 0xDA)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def textbox(slide, l, t, w, h, text, size=18, color=RGBColor(0x33, 0x33, 0x33),
            bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri",
            anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb


def bullets(slide, l, t, w, h, items, size=17, color=RGBColor(0x33, 0x33, 0x33),
            space_after=12, marker="—  "):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = marker + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def kicker(slide, text, color=NAVY):
    textbox(slide, Inches(0.6), Inches(0.35), Inches(9), Inches(0.4),
            text.upper(), size=13, color=color, bold=True)


def title(slide, text, size=28, top=Inches(0.65), color=RGBColor(0x1A, 0x1A, 0x1A)):
    textbox(slide, Inches(0.6), top, Inches(12.1), Inches(0.9), text,
            size=size, color=color, bold=True)


def rule(slide, top=Inches(1.4), color=NAVY, width=Inches(12.1)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, width, Pt(2.2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False


def footer(slide, n):
    textbox(slide, Inches(0.6), Inches(7.12), Inches(8), Inches(0.3),
            "Furto, Renda e Território — UFABC", size=10, color=GREY)
    textbox(slide, Inches(12.0), Inches(7.12), Inches(0.8), Inches(0.3),
            str(n), size=10, color=GREY, align=PP_ALIGN.RIGHT)


def picture_fit(slide, path_or_img, l, t, max_w, max_h):
    """Places an image inside a bounding box, preserving aspect ratio, centered.
    Accepts a file path or a PIL Image."""
    if isinstance(path_or_img, Image.Image):
        im = path_or_img
        iw, ih = im.size
        buf = BytesIO()
        im.save(buf, format="PNG")
        buf.seek(0)
        src = buf
    else:
        with Image.open(path_or_img) as im:
            iw, ih = im.size
        src = path_or_img
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = int(w / ar)
    else:
        h = max_h
        w = int(h * ar)
    x = l + int((max_w - w) / 2)
    y = t + int((max_h - h) / 2)
    slide.shapes.add_picture(src, x, y, width=w, height=h)


def img(name):
    return os.path.join(OUT_DIR, name)


def crop_top(name, frac):
    """Crops the top `frac` of an outputs/ image (used to isolate a panel row)."""
    im = Image.open(img(name))
    w, h = im.size
    return im.crop((0, 0, w, int(h * frac)))


def stat_pill(slide, l, t, w, h, value, label, value_color=NAVY, value_size=26):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    box.adjustments[0] = 0.08
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GREY
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(value_size)
    r1.font.bold = True
    r1.font.color.rgb = value_color
    r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(12)
    r2.font.color.rgb = GREY
    r2.font.name = "Calibri"


def simple_table(slide, l, t, w, h, headers, rows, col_widths=None,
                  header_size=13, cell_size=13, align_center=False,
                  highlight_row=None):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), l, t, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for c, cw in zip(tbl.columns, col_widths):
            c.width = cw
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = htext
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]
        if align_center:
            p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(header_size)
    for i, row in enumerate(rows, start=1):
        is_hl = highlight_row is not None and (i - 1) == highlight_row
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN_SOFT if is_hl else (WHITE if i % 2 else LIGHT_GREY)
            p = cell.text_frame.paragraphs[0]
            if align_center:
                p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(cell_size)
            p.runs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            p.runs[0].font.bold = is_hl
    return tbl_shape


# =====================================================================
# 1. TÍTULO
# =====================================================================
s = add_slide(); set_bg(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background(); band.shadow.inherit = False

textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.4),
        "ANÁLISE DE DADOS PARA PLANEJAMENTO TERRITORIAL • UFABC", size=13, color=GREY, bold=True)
textbox(s, Inches(0.9), Inches(2.65), Inches(11.5), Inches(1.1),
        "Furto, Renda e Território", size=44, color=NAVY, bold=True)
textbox(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6),
        "Uma análise espacial da criminalidade patrimonial nos municípios de São Paulo",
        size=20, color=RGBColor(0x33, 0x33, 0x33))
textbox(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.5),
        "Equipe: João Vitor Oliveira, Julio Cesar Salvino, Lucas Amorim", size=15, color=GREY)
textbox(s, Inches(0.9), Inches(5.25), Inches(11.5), Inches(0.5),
        "Orientação: Prof.ª Flávia da Fonseca Feitosa", size=15, color=GREY)

# =====================================================================
# 2. MOTIVAÇÃO
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Hipótese e motivação")
title(s, "Motivação: a pergunta que originou o trabalho")
rule(s)
textbox(s, Inches(0.6), Inches(1.75), Inches(6.7), Inches(0.9),
        '"Acontecem mais crimes nas\ncidades mais pobres?"', size=24, bold=True,
        color=RGBColor(0x1A, 0x1A, 0x1A))
textbox(s, Inches(0.6), Inches(2.75), Inches(6.7), Inches(0.7),
        "Hipótese inicial (H1): cidades mais pobres registram mais crimes patrimoniais.",
        size=16, color=NAVY, bold=True)
bullets(s, Inches(0.6), Inches(3.55), Inches(6.7), Inches(1.8), [
    "Ocorrências criminais — SSP-SP, 2025 (roubo, furto e lesão corporal).",
    "Renda média per capita — Censo IBGE 2022.",
    "645 municípios — todo o estado de São Paulo.",
], size=16, space_after=10)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.5), Inches(6.7), Inches(1.35))
box.adjustments[0] = 0.08
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFD, 0xEE, 0xE3); box.line.color.rgb = ORANGE; box.line.width = Pt(1)
box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; r = p.add_run()
r.text = ("Observação inicial: a taxa de furto (por 1.000 hab.) mostra concentração "
          "aparente na macrometrópole e no litoral — não no interior mais pobre.")
r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x66, 0x33, 0x00); r.font.name = "Calibri"
picture_fit(s, img("mapa_taxa_furto.png"), Inches(7.7), Inches(1.75), Inches(4.9), Inches(5.1))
footer(s, 2)

# =====================================================================
# 3. DADOS E VARIÁVEIS
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Dados e variáveis")
title(s, "Dados e variáveis")
rule(s)
rows = [
    ("Taxa de roubo (art. 157)", "SSP-SP, 2025", "Ocorrências / 1.000 hab."),
    ("Taxa de furto (art. 155)", "SSP-SP, 2025", "Ocorrências / 1.000 hab. — variável dependente principal"),
    ("Taxa de lesão corporal (art. 129)", "SSP-SP, 2025", "Ocorrências / 1.000 hab."),
    ("Renda média per capita", "IBGE, Censo 2022", "R$ por habitante"),
    ("População", "IBGE, Tabela 4709", "Denominador das taxas"),
    ("Litorâneo / turístico", "Classificação do grupo", "Dummy (1/0) — 16 municípios da costa"),
]
simple_table(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(3.6),
             ["Variável", "Fonte", "Descrição"], rows,
             col_widths=[Inches(3.4), Inches(2.4), Inches(6.3)], header_size=14, cell_size=13)
bullets(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.2), [
    "Chave de junção: código IBGE de 7 dígitos — nunca o nome (a SSP abrevia, "
    "ex. \"S BARBARA D OESTE\"), o que perderia observações na junção.",
], size=15, color=GREY)
footer(s, 3)

# =====================================================================
# 4. POR QUE O FURTO — OUTLIERS POR CRIME
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Dados e variáveis")
title(s, "Por que o furto: cada crime tem seus outliers")
rule(s)
picture_fit(s, crop_top("dispersoes_taxas_renda.png", 0.545), Inches(0.5), Inches(1.7), Inches(12.3), Inches(3.5))
bullets(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.9), [
    "Roubo: o outlier mais extremo é de um município não-litorâneo; os litorâneos "
    "(laranja) aparecem elevados, mas não dominam o topo.",
    "Furto: os outliers extremos são quase todos litorâneos — Mongaguá e outros "
    "municípios da costa aparecem bem acima da nuvem geral de pontos.",
    "Lesão corporal: o outlier extremo é de renda baixíssima e não-litorâneo; os "
    "litorâneos ficam perdidos no meio da distribuição.",
], size=14, space_after=8)
footer(s, 4)

# =====================================================================
# 5. CORRELAÇÃO
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Dados e variáveis")
title(s, "A renda correlaciona positivamente com o furto")
rule(s)
crows = [
    ("Roubo × Renda", "0,302", "0,389"),
    ("Furto × Renda", "0,359", "0,380"),
    ("Lesão corporal × Renda", "−0,190", "−0,166"),
]
simple_table(s, Inches(0.6), Inches(1.9), Inches(6.4), Inches(1.9),
             ["Par", "Pearson", "Spearman"], crows,
             col_widths=[Inches(3.2), Inches(1.6), Inches(1.6)],
             header_size=14, cell_size=14, align_center=False)
textbox(s, Inches(0.6), Inches(4.1), Inches(6.4), Inches(1.0),
        "Roubo e furto sobem junto com a renda — o oposto do que H1 previa. "
        "Só a lesão corporal cai, e fraco.", size=16, color=RGBColor(0x33, 0x33, 0x33))
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.3), Inches(6.4), Inches(1.5))
box.adjustments[0] = 0.08; box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GREY; box.line.fill.background(); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; r = p.add_run()
r.text = ("Correlação não implica causalidade — mas já contraria H1 para roubo e "
          "furto, e os municípios litorâneos puxam boa parte desse padrão, motivando "
          "a investigação espacial que segue.")
r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = GREY
picture_fit(s, img("correlacao_matriz.png"), Inches(7.3), Inches(1.75), Inches(5.4), Inches(5.1))
footer(s, 5)

# =====================================================================
# 6. AUTOCORRELAÇÃO ESPACIAL
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Desenvolvimento e resultados")
title(s, "O crime tem endereço: autocorrelação espacial")
rule(s)
picture_fit(s, img("mapa_lisa_furto.png"), Inches(0.6), Inches(1.75), Inches(6.9), Inches(5.1))
stat_pill(s, Inches(7.8), Inches(1.8), Inches(4.8), Inches(1.0), "I = 0,832", "Moran's I — taxa de roubo (forte)")
stat_pill(s, Inches(7.8), Inches(2.95), Inches(4.8), Inches(1.0), "I = 0,422", "Moran's I — taxa de furto (moderada)")
stat_pill(s, Inches(7.8), Inches(4.1), Inches(4.8), Inches(1.0), "I = 0,181", "Moran's I — lesão corporal (fraca)")
textbox(s, Inches(7.8), Inches(5.3), Inches(4.8), Inches(1.6),
        "Todos significativos a p = 0,001 (999 permutações): o crime não está "
        "distribuído ao acaso. LISA (furto): 36 municípios em cluster Alto-Alto "
        "acompanham o litoral (36% já são litorâneos/turísticos) — primeira pista "
        "quantitativa da H2.",
        size=13, color=GREY)
footer(s, 6)

# =====================================================================
# 7. REGRESSÃO CLÁSSICA
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Desenvolvimento e resultados")
title(s, "Regressão clássica: a renda sozinha não basta")
rule(s)
picture_fit(s, img("reg_multipla.png"), Inches(0.6), Inches(1.75), Inches(7.0), Inches(5.1))
bullets(s, Inches(7.9), Inches(1.85), Inches(4.7), Inches(3.0), [
    "Simples: taxa = 0,971 + 0,00247·renda   R² = 0,129",
    "Com litoral: taxa = 1,036 + 0,00233·renda + 10,97·litoral   R² = 0,357",
    "O litoral soma ~11 furtos/1.000 hab., equivalente a +R$4.700 de renda.",
], size=14, space_after=12)
textbox(s, Inches(7.9), Inches(5.0), Inches(4.7), Inches(1.9),
        "Distância de Cook: 29 municípios influentes, 8 deles litorâneos — 50% dos "
        "16 municípios do litoral. O padrão litorâneo não é efeito de poucos "
        "pontos: é estrutural.",
        size=13, color=GREY)
footer(s, 7)

# =====================================================================
# 8. DO GLOBAL AO LOCAL (compacto: spatial error / GWR / GWR+litoral)
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Desenvolvimento e resultados")
title(s, "O espaço importa: do global ao local")
rule(s)
textbox(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.6),
        "Os resíduos do modelo clássico têm Moran's I = 0,38 (p<0,001): a "
        "dependência espacial é ignorada pelo OLS. Como resolvemos isso?",
        size=15, color=RGBColor(0x33, 0x33, 0x33))

def box_model(slide, l, t, w, h, titulo, linhas, r2_val, destaque=False):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = NAVY if destaque else WHITE
    box.line.color.rgb = NAVY if not destaque else NAVY
    box.line.width = Pt(1) if not destaque else Pt(0)
    box.shadow.inherit = False
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = ORANGE if not destaque else WHITE
    top_bar.line.fill.background(); top_bar.shadow.inherit = False
    tcol = WHITE if destaque else RGBColor(0x1A, 0x1A, 0x1A)
    lcol = RGBColor(0xDD, 0xE6, 0xF0) if destaque else GREY
    textbox(slide, l + Inches(0.25), t + Inches(0.2), w - Inches(0.5), Inches(0.4),
            titulo, size=15, bold=True, color=tcol)
    textbox(slide, l + Inches(0.25), t + Inches(0.65), w - Inches(0.5), h - Inches(1.5),
            "\n".join(linhas), size=12.5, color=lcol)
    textbox(slide, l + Inches(0.25), t + h - Inches(0.75), w - Inches(0.5), Inches(0.6),
            r2_val, size=22, bold=True, color=WHITE if destaque else ORANGE, align=PP_ALIGN.RIGHT)

box_model(s, Inches(0.6), Inches(2.55), Inches(3.9), Inches(2.0),
          "1. Spatial Error (global)",
          ["Testes de Lagrange apontam este", "modelo como mais adequado.",
           "λ = 0,588 — resíduo cai a −0,03"],
          "R² = 0,392")
box_model(s, Inches(4.7), Inches(2.55), Inches(3.9), Inches(2.0),
          "2. GWR (renda)",
          ["Banda adaptativa (~10 vizinhos).", "O coeficiente da renda passa a",
           "variar livremente no espaço."],
          "R² = 0,458")
box_model(s, Inches(8.8), Inches(2.55), Inches(3.9), Inches(2.0),
          "3. GWR + litoral",
          ["Banda adaptativa incorporando", "renda e a condição litorânea/", "turística, localmente."],
          "R² = 0,532", destaque=True)
textbox(s, Inches(0.6), Inches(5.05), Inches(12.1), Inches(0.6),
        "Melhor ajuste entre todos os modelos testados — e o menor AIC (3027,4).",
        size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
picture_fit(s, img("reg_esp_residuos_ols.png"), Inches(4.7), Inches(5.85), Inches(3.9), Inches(1.15))
textbox(s, Inches(4.7), Inches(5.72), Inches(3.9), Inches(0.25),
        "resíduos do OLS: por que sair dele", size=10.5, italic=True, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 8)

# =====================================================================
# 9. GWR + LITORAL — CONTRAFACTUAL
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Desenvolvimento e resultados")
title(s, "GWR + litoral: isolando o efeito turístico")
rule(s)
picture_fit(s, img("gwr_contrafactual_dif.png"), Inches(0.6), Inches(1.75), Inches(6.4), Inches(5.1))
textbox(s, Inches(7.3), Inches(1.9), Inches(5.4), Inches(0.6),
        "Cenário contrafactual", size=17, bold=True, color=NAVY)
bullets(s, Inches(7.3), Inches(2.5), Inches(5.4), Inches(2.0), [
    "O que acontece se \"zerarmos\" a variável litoral no GWR e recalcularmos a predição?",
    "A taxa cai em até 12 furtos/1.000 hab. — impacto concentrado inteiramente na faixa costeira.",
], size=14, space_after=12)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(4.7), Inches(5.4), Inches(2.0))
box.adjustments[0] = 0.08; box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GREY; box.line.fill.background(); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Resíduo médio sem a variável litoral:"
r.font.size = Pt(13); r.font.color.rgb = GREY
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = "+10,68 nos municípios litorâneos    −0,27 nos demais"
r2.font.size = Pt(14); r2.font.bold = True; r2.font.color.rgb = ORANGE
p3 = tf.add_paragraph(); r3 = p3.add_run()
r3.text = "Com a variável incluída, ambos os grupos centram em zero."
r3.font.size = Pt(13); r3.font.color.rgb = GREY
footer(s, 9)

# =====================================================================
# 10. MECANISMO
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Desenvolvimento e resultados")
title(s, "Onde o modelo erra: o rastro da população flutuante")
rule(s)
picture_fit(s, img("residuos_mapa_sem_litoral.png"), Inches(0.6), Inches(1.75), Inches(6.4), Inches(5.1))
textbox(s, Inches(7.3), Inches(1.9), Inches(5.4), Inches(0.9),
        "O modelo (sem litoral) subestima as ocorrências exatamente no litoral — a "
        "faixa costeira concentra a subestimação mais intensa do estado.",
        size=14, color=RGBColor(0x33, 0x33, 0x33))
textbox(s, Inches(7.3), Inches(3.0), Inches(5.4), Inches(0.4),
        "O mecanismo numérico", size=16, bold=True, color=NAVY)
bullets(s, Inches(7.3), Inches(3.5), Inches(5.4), Inches(1.8), [
    "Numerador (crimes): as ocorrências inflam com o fluxo de turistas no verão.",
    "Denominador (população): o modelo conta só a população residente medida pelo Censo IBGE.",
], size=14, space_after=10)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(5.5), Inches(5.4), Inches(1.2))
box.adjustments[0] = 0.08; box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFD, 0xEE, 0xE3); box.line.color.rgb = ORANGE; box.line.width = Pt(1); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; r = p.add_run()
r.text = "A taxa por 1.000 hab. superestima o risco real para quem mora lá — sem o dummy, o modelo sempre vai subestimar a taxa observada."
r.font.size = Pt(13); r.font.color.rgb = RGBColor(0x66, 0x33, 0x00)
footer(s, 10)

# =====================================================================
# 11. EXEMPLO: SÃO PAULO x MONGAGUÁ
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Desenvolvimento e resultados")
title(s, "Quanto R$ 1 de renda pesa — e onde falha")
rule(s)
erows = [
    ("Renda", "R$ 4.547,62", "R$ 2.138,21"),
    ("Banda (GWR)", "23,6 km", "41,9 km"),
    ("β renda", "0,00156", "0,00019"),
    ("t (signif.)", "3,28 (significativo)", "0,46 (não significativo)"),
    ("R² local", "0,31", "0,18"),
    ("Furto previsto", "13,7", "13,2"),
    ("Resíduo", "+11,5", "+21,0"),
]
simple_table(s, Inches(0.6), Inches(1.85), Inches(6.6), Inches(4.6),
             ["Métrica", "São Paulo", "Mongaguá"], erows,
             col_widths=[Inches(2.2), Inches(2.2), Inches(2.2)],
             header_size=14, cell_size=13, align_center=True)
bullets(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(4.2), [
    "Em São Paulo e Mongaguá, a previsão do furto é parecida (~13 a 14) — mas só "
    "em São Paulo o coeficiente da renda é confiável (t=3,28).",
    "Em Mongaguá, t=0,46 diz que a renda local não explica o crime. O resíduo de "
    "+21 é quase todo população turística flutuante.",
    "Deixar a renda variar no espaço faz o R² saltar de 0,13 (global) para 0,46 "
    "(GWR local): não é só \"quanto\" de renda existe, é \"onde\" ela está.",
], size=14, space_after=14)
footer(s, 11)

# =====================================================================
# 12. COMPARAÇÃO FINAL DOS MODELOS
# =====================================================================
s = add_slide(); set_bg(s)
kicker(s, "Desenvolvimento e resultados")
title(s, "Comparação final dos modelos")
rule(s)
rows = [
    ("OLS", "0,129", "0,357", "+0,228"),
    ("Spatial Error", "0,392", "0,444", "+0,052"),
    ("GWR", "0,458", "0,532", "+0,074"),
]
tbl_shape = s.shapes.add_table(len(rows) + 1, 4, Inches(1.9), Inches(2.2), Inches(9.5), Inches(2.2))
tbl = tbl_shape.table
for c, w in zip(tbl.columns, [Inches(3.3), Inches(2.2), Inches(2.2), Inches(1.8)]):
    c.width = w
headers = ["Modelo", "Sem litoral", "Com litoral", "Δ"]
for j, htext in enumerate(headers):
    c = tbl.cell(0, j); c.text = htext
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.color.rgb = WHITE; p.runs[0].font.bold = True; p.runs[0].font.size = Pt(16)
for i, row in enumerate(rows, start=1):
    is_best = row[0] == "GWR"
    for j, val in enumerate(row):
        cell = tbl.cell(i, j); cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN_SOFT if is_best else (WHITE if i % 2 else LIGHT_GREY)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(17); p.runs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        p.runs[0].font.bold = is_best
textbox(s, Inches(1.9), Inches(4.9), Inches(9.5), Inches(1.4),
        "O ganho do litoral é maior no modelo mais simples (OLS, +0,228) e menor nos "
        "modelos espaciais (+0,052 no spatial error, +0,074 no GWR) — sinal de que "
        "eles já capturavam parte do padrão litorâneo implicitamente, pela geografia.",
        size=15, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 12)

# =====================================================================
# 13. CONCLUSÃO
# =====================================================================
s = add_slide(); set_bg(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.32), SW, Inches(0.18))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background(); band.shadow.inherit = False
kicker(s, "Conclusão")
title(s, "Conclusão: não é pobreza. É território.")
rule(s)

pontos = [
    "A hipótese inicial (H1) não se sustenta: a relação entre renda e furto é "
    "positiva em todos os modelos testados, não negativa.",
    "A renda sozinha explica pouco (R²=0,13) e deixa forte autocorrelação espacial "
    "nos resíduos clássicos.",
    "Essa estrutura é majoritariamente explicada pela condição litorânea/turística "
    "— a população flutuante infla a taxa sobre um denominador de residentes fixos.",
    "A própria relação renda–furto não é estacionária: forte no interior, quase "
    "nula no litoral, como o GWR demonstra.",
    "Modelo vencedor: GWR com renda e litoral, R²=0,532 — o melhor ajuste entre "
    "todos os testados.",
]
top = Inches(1.85)
for i, txt in enumerate(pontos):
    y = top + Inches(0.92) * i
    is_last = i == len(pontos) - 1
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.55), Inches(0.55))
    num.fill.solid(); num.fill.fore_color.rgb = GREEN if is_last else NAVY
    num.line.fill.background(); num.shadow.inherit = False
    tf = num.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run(); r.text = str(i + 1)
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
    textbox(s, Inches(1.35), y - Inches(0.02), Inches(11.4), Inches(0.85), txt,
            size=14.5, color=RGBColor(0x1A, 0x1A, 0x1A) if not is_last else GREEN,
            bold=is_last)
footer(s, 13)

out_path = os.path.join(SLIDES_DIR, "apresentacao_final.pptx")
prs.save(out_path)
print(f"OK: {len(prs.slides)} slides")
print(f"saved: {out_path}")
